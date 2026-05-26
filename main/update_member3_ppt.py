#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将数值实验与政策仿真内容追加到 report/模型.pptx。

运行方式：
  /opt/anaconda3/bin/python3 main/update_member3_ppt.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "report" / "模型.pptx"
IMG_DIR = ROOT / "main"

BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(40, 49, 59)
MUTED = RGBColor(95, 108, 124)
ACCENT = RGBColor(208, 97, 44)
LIGHT_BLUE = RGBColor(232, 239, 248)

MEMBER3_TITLES = {
    "五、数值实验与政策仿真",
    "实验设计与评价指标",
    "需求不确定性影响",
    "补贴灵敏度分析",
    "CV扩展与管理洞察",
    "管理洞察与政策建议",
}


def slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            texts.append(shape.text.strip())
    return "\n".join(t for t in texts if t)


def add_textbox(slide, left, top, width, height, text, font_size=20,
                color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def set_template_title(slide, title):
    title_shape = None
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and shape.top < Inches(0.6) and shape.left > Inches(3.5):
            title_shape = shape
            break
    if title_shape is None:
        title_shape = slide.shapes[0]
    title_shape.text = title
    frame = title_shape.text_frame
    frame.word_wrap = True
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = BLUE


def remove_content_placeholder(slide):
    for shape in list(slide.shapes):
        if (
            getattr(shape, "is_placeholder", False)
            and shape.top >= Inches(0.8)
            and shape.width >= Inches(7.0)
        ):
            shape._element.getparent().remove(shape._element)


def add_subtitle(slide, subtitle):
    add_textbox(slide, Inches(0.72), Inches(0.86), Inches(7.15), Inches(0.28),
                subtitle, font_size=12, color=MUTED)


def add_bullets(slide, left, top, width, height, bullets, font_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
    return box


def add_metric_row(slide, top, metrics):
    start_left = Inches(0.72)
    gap = Inches(0.14)
    width = Inches(1.76)
    for i, (label, value, color) in enumerate(metrics):
        left = start_left + i * (width + gap)
        shape = slide.shapes.add_shape(1, left, top, width, Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = RGBColor(205, 216, 230)
        add_textbox(slide, left + Inches(0.12), top + Inches(0.10), width - Inches(0.24), Inches(0.20),
                    label, font_size=10, color=MUTED, bold=False, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.12), top + Inches(0.32), width - Inches(0.24), Inches(0.28),
                    value, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)


def remove_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    prs.slides._sldIdLst.remove(slide_id)


def add_member3_slides(prs):
    section_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    new_slides = []

    slide = prs.slides.add_slide(section_layout)
    slide.shapes[0].text = "五、数值实验与政策仿真"
    for paragraph in slide.shapes[0].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = BLUE
    new_slides.append(slide)

    slide = prs.slides.add_slide(content_layout)
    set_template_title(slide, "实验设计与评价指标")
    remove_content_placeholder(slide)
    add_subtitle(slide, "围绕随机规划解的投资响应与政策敏感性展开")
    add_bullets(slide, Inches(0.72), Inches(1.18), Inches(5.2), Inches(2.05), [
        "实验一：比较随机规划 DEP 与均值模型 EV，评估需求波动是否促使企业预留更多 EYC。",
        "实验二：改变补贴比例 k，观察企业成本、政府预算、改造区域数与总排放的响应。",
        "实验三：设置不同需求波动水平 CV，分析补贴边际效用是否随不确定性增强而变化。",
    ], font_size=16)
    add_metric_row(slide, Inches(3.72), [
        ("核心对比", "DEP vs EV", BLUE),
        ("补贴网格", "k=0~1", ACCENT),
        ("波动水平", "CV=0.05/0.15/0.30", RGBColor(62, 138, 94)),
        ("主算例", "3区×24期×20场景", RGBColor(124, 90, 166)),
    ])
    add_bullets(slide, Inches(6.28), Inches(1.28), Inches(2.35), Inches(2.25), [
        "评价指标：EYC 投入、VSS、期望总排放、政府补贴支出。",
        "所有实验使用 Gurobi 直接求解确定性等价模型。",
        "图表展示全网格趋势，表格保留关键节点。",
    ], font_size=12.8)
    new_slides.append(slide)

    slide = prs.slides.add_slide(content_layout)
    set_template_title(slide, "需求不确定性影响")
    remove_content_placeholder(slide)
    add_subtitle(slide, "随机规划 DEP 与均值模型 EV 的一阶段投资对比")
    slide.shapes.add_picture(str(IMG_DIR / "member3_uncertainty_impact.png"),
                             Inches(0.72), Inches(1.22), width=Inches(5.78))
    add_bullets(slide, Inches(6.65), Inches(1.30), Inches(2.15), Inches(2.55), [
        "结论：平均需求模型系统性低估高峰风险。",
        "中规模 EYC 投入从 11.20 台提高到 16.19 台。",
        "大规模缓冲增量达到 7.32 台，风险缓冲随规模扩大。",
        "VSS 为正，说明随机规划带来可量化经济价值。",
    ], font_size=12.5)
    add_metric_row(slide, Inches(4.18), [
        ("小规模增量", "+3.28台", BLUE),
        ("中规模增量", "+4.99台", BLUE),
        ("大规模增量", "+7.32台", BLUE),
        ("中规模VSS", "29258.29", ACCENT),
    ])
    new_slides.append(slide)

    slide = prs.slides.add_slide(content_layout)
    set_template_title(slide, "补贴灵敏度分析")
    remove_content_placeholder(slide)
    add_subtitle(slide, "中规模实例：k=0,0.1,...,1.0")
    slide.shapes.add_picture(str(IMG_DIR / "member3_subsidy_sensitivity.png"),
                             Inches(0.70), Inches(1.12), width=Inches(5.95))
    add_bullets(slide, Inches(6.82), Inches(1.28), Inches(2.05), Inches(3.15), [
        "结论：补贴降低企业负担，但不是充分减排条件。",
        "k=0 到 k=0.9 时，EYC 投入和总排放基本不变。",
        "电网改造区域数稳定为 3 个，说明基准方案已跨过改造门槛。",
        "k=1.0 时 EYC 投入小幅上升到 17.20 台。",
        "按 0.5% 减排标准，当前参数下未形成显著政策临界点。",
    ], font_size=11.6)
    new_slides.append(slide)

    slide = prs.slides.add_slide(content_layout)
    set_template_title(slide, "CV扩展与管理洞察")
    remove_content_placeholder(slide)
    add_subtitle(slide, "不同需求波动水平下的补贴边际效用")
    slide.shapes.add_picture(str(IMG_DIR / "member3_cv_marginal_effect.png"),
                             Inches(0.72), Inches(1.22), width=Inches(5.78))
    add_bullets(slide, Inches(6.65), Inches(1.28), Inches(2.15), Inches(2.95), [
        "结论：不确定性越高，EYC 更像“保险性产能”。",
        "CV=0.30 时，k=0 的 EYC 投入已达到 22.74 台。",
        "补贴可扩大电动化投入，但未明显改变排放曲线。",
        "减排效果取决于设备是否真正替代 DYC 作业。",
    ], font_size=12.5)
    add_metric_row(slide, Inches(4.18), [
        ("CV=0.05 EYC", "12.63台", BLUE),
        ("CV=0.15 EYC", "16.19台", BLUE),
        ("CV=0.30 EYC", "22.74台", BLUE),
        ("结论", "补贴需协同碳政策", ACCENT),
    ])
    new_slides.append(slide)

    slide = prs.slides.add_slide(content_layout)
    set_template_title(slide, "管理洞察与政策建议")
    remove_content_placeholder(slide)
    add_subtitle(slide, "从“补贴多少”转向“补贴如何触发真实减排”")

    add_bullets(slide, Inches(0.72), Inches(1.18), Inches(3.95), Inches(3.25), [
        "企业视角：需求波动越强，越需要把 EYC 作为高峰风险缓冲，而不是只按平均需求配置。",
        "政府视角：投资补贴能降低企业成本，但若电动化已具备经济性，继续提高补贴的边际减排有限。",
        "政策设计：补贴应与运营侧约束联动，例如 EYC 作业占比、实际减排吨数、碳价或免费配额递减。",
    ], font_size=13.5)

    add_textbox(slide, Inches(5.12), Inches(1.22), Inches(3.55), Inches(0.32),
                "政策工具的作用边界", font_size=15, color=BLUE, bold=True)
    policy_rows = [
        ("投资补贴 k", "降低前期成本", "对已改造区域边际减排弱"),
        ("碳交易价格", "提高DYC作业成本", "直接影响运营排放结构"),
        ("配额递减", "压缩高排放空间", "推动持续减排"),
        ("绩效型补贴", "绑定实际减排", "避免只补贴存量投资"),
    ]
    top = Inches(1.72)
    row_h = Inches(0.58)
    for idx, (tool, effect, insight) in enumerate(policy_rows):
        y = top + idx * row_h
        bg = LIGHT_BLUE if idx % 2 == 0 else RGBColor(246, 248, 251)
        rect = slide.shapes.add_shape(1, Inches(5.10), y, Inches(3.72), row_h - Inches(0.04))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(218, 226, 236)
        add_textbox(slide, Inches(5.22), y + Inches(0.09), Inches(0.88), Inches(0.22),
                    tool, font_size=9.5, color=BLUE, bold=True)
        add_textbox(slide, Inches(6.15), y + Inches(0.09), Inches(1.05), Inches(0.22),
                    effect, font_size=9.5, color=DARK)
        add_textbox(slide, Inches(7.18), y + Inches(0.09), Inches(1.48), Inches(0.22),
                    insight, font_size=9.2, color=DARK)

    add_metric_row(slide, Inches(4.38), [
        ("核心发现", "随机规划更稳健", BLUE),
        ("补贴定位", "成本分担", ACCENT),
        ("减排关键", "运营替代", RGBColor(62, 138, 94)),
        ("建议", "政策组合", RGBColor(124, 90, 166)),
    ])
    new_slides.append(slide)

    return new_slides


def move_slides_before_appendix(prs, new_slides):
    appendix_index = len(prs.slides) - len(new_slides)
    for i, slide in enumerate(prs.slides):
        if "附录和参考文献" in slide_text(slide):
            appendix_index = i
            break

    sld_id_list = prs.slides._sldIdLst
    sld_id_elements = list(sld_id_list)
    new_ids = sld_id_elements[-len(new_slides):]
    for sld_id in new_ids:
        sld_id_list.remove(sld_id)
    for offset, sld_id in enumerate(new_ids):
        sld_id_list.insert(appendix_index + offset, sld_id)


def main():
    prs = Presentation(PPT_PATH)
    for index in reversed(range(len(prs.slides))):
        text = slide_text(prs.slides[index])
        if any(title in text for title in MEMBER3_TITLES):
            remove_slide(prs, index)

    new_slides = add_member3_slides(prs)
    move_slides_before_appendix(prs, new_slides)
    prs.save(PPT_PATH)
    print(f"已更新: {PPT_PATH}")
    print(f"新增页面: {len(new_slides)}")


if __name__ == "__main__":
    main()
